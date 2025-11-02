import streamlit as st
from openai import OpenAI
import PyPDF2
from docx import Document
import openpyxl
import io
from pptx import Presentation
import json

st.title("📄 ドキュメント質問応答 & 🤖 チャットボット")
st.write(
    "ドキュメントQA（ファイル質問）とチャットボット（自由会話）の両方が使えます。"
    "このアプリを利用するには OpenAI API キーが必要です。取得方法は[こちら](https://platform.openai.com/account/api-keys)。"
)

openai_api_key = st.text_input("OpenAI APIキー", type="password")
if not openai_api_key:
    st.info("OpenAI APIキーを入力してください。", icon="🗝️")
else:
    client = OpenAI(api_key=openai_api_key)

    # タブ切り替え
    tab1, tab2 = st.tabs(["ドキュメントQA", "チャットボット"])

    # 履歴データの初期化
    if "history_doc" not in st.session_state:
        st.session_state["history_doc"] = []
    if "history_chat" not in st.session_state:
        st.session_state["history_chat"] = []

    # -------- ドキュメントQAタブ -------- #
    with tab1:
        uploaded_file = st.file_uploader(
            "ドキュメントをアップロードしてください（.txt, .md, .pdf, .docx, .xlsx, .pptx）", 
            type=("txt", "md", "pdf", "docx", "xlsx", "pptx"),
            key="doc_uploader"
        )
        question = st.text_area(
            "ドキュメントについて質問してください！",
            placeholder="この文書の要約を教えてください。",
            disabled=not uploaded_file,
            key="doc_question"
        )

        def extract_text(file):
            filename = file.name
            ext = filename.split('.')[-1].lower()
            if ext in ['txt', 'md']:
                return file.read().decode()
            elif ext == 'pdf':
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() or ""
                return text
            elif ext == 'docx':
                doc = Document(io.BytesIO(file.read()))
                return "\n".join([para.text for para in doc.paragraphs])
            elif ext == 'xlsx':
                wb = openpyxl.load_workbook(io.BytesIO(file.read()), data_only=True)
                text = ""
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        text += " ".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
                return text
            elif ext == 'pptx':
                prs = Presentation(io.BytesIO(file.read()))
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            else:
                return None

        if uploaded_file and question:
            document = extract_text(uploaded_file)
            if not document or document.strip() == "":
                st.error("ファイルからテキストを抽出できませんでした。")
            else:
                messages = [
                    {
                        "role": "user",
                        "content": f"以下はドキュメントです: {document} \n\n---\n\n {question}",
                    }
                ]
                stream = client.chat.completions.create(
                    model="gpt-3.5-turbo",
                    messages=messages,
                    stream=True,
                )
                answer = st.write_stream(stream)
                st.session_state["history_doc"].append({
                    "question": question,
                    "answer": answer,
                    "bookmark": False,
                    "doc_name": uploaded_file.name
                })

        st.header("履歴（ドキュメントQA）")
        for i, h in enumerate(st.session_state["history_doc"]):
            col1, col2 = st.columns([10, 1])
            with col1:
                st.write(f"**Q:** {h['question']}")
                st.write(f"**A:** {h['answer']}")
                st.write(f"**ドキュメント:** {h['doc_name']}")
            with col2:
                if st.button("⭐" if h["bookmark"] else "☆", key=f"bookmark_doc_{i}"):
                    st.session_state["history_doc"][i]["bookmark"] = not h["bookmark"]

        st.header("しおり一覧（ドキュメントQA）")
        for h in [x for x in st.session_state["history_doc"] if x["bookmark"]]:
            st.write(f"**Q:** {h['question']}")
            st.write(f"**A:** {h['answer']}")
            st.write(f"**ドキュメント:** {h['doc_name']}")

        history_json = json.dumps(st.session_state["history_doc"], ensure_ascii=False, indent=2)
        st.download_button("履歴をダウンロード（ドキュメントQA）", data=history_json, file_name="history_doc.json", mime="application/json")

    # -------- チャットボットタブ -------- #
    with tab2:
        user_message = st.text_area(
            "チャットを入力してください",
            placeholder="自由に質問や会話をしてください。",
            key="chat_message"
        )

        if user_message:
            # 直近の履歴を使って会話
            messages = [{"role": "user", "content": user_message}]
            stream = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=messages,
                stream=True,
            )
            answer = st.write_stream(stream)
            st.session_state["history_chat"].append({
                "question": user_message,
                "answer": answer,
                "bookmark": False,
            })

        st.header("履歴（チャットボット）")
        for i, h in enumerate(st.session_state["history_chat"]):
            col1, col2 = st.columns([10, 1])
            with col1:
                st.write(f"**Q:** {h['question']}")
                st.write(f"**A:** {h['answer']}")
            with col2:
                if st.button("⭐" if h["bookmark"] else "☆", key=f"bookmark_chat_{i}"):
                    st.session_state["history_chat"][i]["bookmark"] = not h["bookmark"]

        st.header("しおり一覧（チャットボット）")
        for h in [x for x in st.session_state["history_chat"] if x["bookmark"]]:
            st.write(f"**Q:** {h['question']}")
            st.write(f"**A:** {h['answer']}")

        history_json_chat = json.dumps(st.session_state["history_chat"], ensure_ascii=False, indent=2)
        st.download_button("履歴をダウンロード（チャットボット）", data=history_json_chat, file_name="history_chat.json", mime="application/json")
